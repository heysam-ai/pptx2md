from __future__ import print_function

import os
import re
from operator import attrgetter
from typing import Optional

from PIL import Image
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from rapidfuzz import process as fuze_process
from tqdm import tqdm

from pptx2md.global_var import g

picture_count = 0

global out

LINK_PATTERN = re.compile(r"(https?://|www\.)\S+")


# pptx type defination rules
def is_title(shape):
    if shape.is_placeholder and (
        shape.placeholder_format.type == PP_PLACEHOLDER.TITLE
        or shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE
        or shape.placeholder_format.type == PP_PLACEHOLDER.VERTICAL_TITLE
        or shape.placeholder_format.type == PP_PLACEHOLDER.CENTER_TITLE
    ):
        return True
    return False


def is_text_block(shape):
    if shape.has_text_frame:
        if (
            shape.is_placeholder
            and shape.placeholder_format.type == PP_PLACEHOLDER.BODY
        ):
            return True
        if len(shape.text) > g.text_block_threshold:
            return True
    return False


def is_list_block(shape):
    levels = []
    for para in shape.text_frame.paragraphs:
        if para.level not in levels:
            levels.append(para.level)
        if para.level != 0 or len(levels) > 1:
            return True
    return False


def is_accent(font):
    if (
        font.underline
        or font.italic
        or (
            font.color.type == MSO_COLOR_TYPE.SCHEME
            and (
                font.color.theme_color == MSO_THEME_COLOR.ACCENT_1
                or font.color.theme_color == MSO_THEME_COLOR.ACCENT_2
                or font.color.theme_color == MSO_THEME_COLOR.ACCENT_3
                or font.color.theme_color == MSO_THEME_COLOR.ACCENT_4
                or font.color.theme_color == MSO_THEME_COLOR.ACCENT_5
                or font.color.theme_color == MSO_THEME_COLOR.ACCENT_6
            )
        )
    ):
        return True
    return False


def is_strong(font):
    if font.bold or (
        font.color.type == MSO_COLOR_TYPE.SCHEME
        and (
            font.color.theme_color == MSO_THEME_COLOR.DARK_1
            or font.color.theme_color == MSO_THEME_COLOR.DARK_2
        )
    ):
        return True
    return False


def get_formatted_text(para):
    res = ""
    for run in para.runs:
        text = run.text
        if text == "":
            continue
        if not g.disable_escaping:
            text = out.get_escaped(text)
        try:
            if run.hyperlink.address:
                text = out.get_hyperlink(text, run.hyperlink.address)
        except Exception:
            text = out.get_hyperlink(text, "error:ppt-link-parsing-issue")
        if is_accent(run.font):
            text = out.get_accent(text)
        elif is_strong(run.font):
            text = out.get_strong(text)
        if not g.disable_color:
            if run.font.color.type == MSO_COLOR_TYPE.RGB:
                text = out.get_colored(text, run.font.color.rgb)
        res += text
    return res.strip()


def process_title(shape, slide_idx, level: Optional[int] = None) -> list[str]:
    global out
    notes = []
    text = shape.text_frame.text.strip()
    if not g.disable_newline_removal_from_title:
        text = text.replace("\n", " ")
    if g.use_custom_title:
        res = fuze_process.extractOne(text, g.titles.keys(), score_cutoff=92)
        if not res:
            g.max_custom_title
            out.put_title(text, g.max_custom_title + 1)
        else:
            notes.append(
                f'Title in slide {slide_idx} "{text}" is converted to "{res[0]}" as specified in title file.'
            )
            out.put_title(res[0], g.titles[res[0]])
    else:
        out.put_title(text, level or 1)

    return notes


def process_text_block(shape, _):
    global out
    if is_list_block(shape):
        # generate list block
        for para in shape.text_frame.paragraphs:
            if para.text.strip() == "":
                continue
            text = get_formatted_text(para)
            out.put_list(text, para.level)
        out.write("\n")
    else:
        # generate paragraph block
        for para in shape.text_frame.paragraphs:
            if para.text.strip() == "":
                continue
            text = get_formatted_text(para)
            out.put_para(text)
    return []


def process_notes(text, _):
    global out
    out.put_para("---")
    out.put_para(text)
    return []


def process_picture(shape, slide_idx):
    notes = []
    if g.disable_image:
        return notes
    global picture_count
    global out

    pic_name = g.file_prefix + str(picture_count)
    pic_ext = shape.image.ext
    if not os.path.exists(g.img_path):
        os.makedirs(g.img_path)

    output_path = g.path_name_ext(g.img_path, pic_name, pic_ext)
    common_path = os.path.commonpath([g.out_path, g.img_path])
    img_outputter_path = os.path.relpath(output_path, common_path)
    with open(output_path, "wb") as f:
        f.write(shape.image.blob)
        picture_count += 1

    # normal images
    if pic_ext != "wmf":
        out.put_image(img_outputter_path, g.max_img_width)
        return notes

    # wmf images, try to convert, if failed, output as original
    try:
        Image.open(output_path).save(os.path.splitext(output_path)[0] + ".png")
        out.put_image(os.path.splitext(img_outputter_path)[0] + ".png", g.max_img_width)
        notes.append(f"Image {output_path} in slide {slide_idx} converted to png.")
    except Exception:
        notes.append(
            f"Cannot convert wmf image {output_path} in slide {slide_idx} to png, this probably won't be displayed correctly."
        )
        out.put_image(img_outputter_path, g.max_img_width)
    return notes


def process_table(shape, _):
    global out
    table = [[cell.text for cell in row.cells] for row in shape.table.rows]
    if len(table) > 0:
        out.put_table(table)
    return []


def ungroup_shapes(shapes):
    res = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            res.extend(ungroup_shapes(shape.shapes))
        else:
            res.append(shape)
    return res


# main
def has_too_many_links(text: str) -> bool:
    cleaned_text = re.sub(LINK_PATTERN, "", text)
    ratio: float = count_non_whitespace_chars(
        cleaned_text
    ) / count_non_whitespace_chars(text)
    return ratio < g.has_too_many_links_ratio


def count_non_whitespace_chars(s: str) -> int:
    return sum(1 for c in s if not c.isspace())


def is_section_header(idx, shapes) -> bool:
    if idx == 0:
        # first slide, by convention is not a section header
        return False

    titles = sum(1 for shape in shapes if is_title(shape))
    text_blocks = sum(
        1 for shape in shapes if not is_title(shape) and is_text_block(shape)
    )

    # by convention, we assume that a slide is a section header if it only has one title or text
    # block and has no table. It can have images in it
    return titles + text_blocks == 1 and all(
        shape.shape_type != MSO_SHAPE_TYPE.TABLE for shape in shapes
    )


def parse(prs, outputer):
    global out
    out = outputer
    notes = []
    seen_section_header = False

    for idx, slide in enumerate(tqdm(prs.slides, desc="Converting slides")):
        if g.page is not None and idx + 1 != g.page:
            continue
        shapes = []
        try:
            shapes = sorted(ungroup_shapes(slide.shapes), key=attrgetter("top", "left"))
        except Exception:
            print(
                "Bad shapes encountered in this slide. Please check or move them and try again."
            )
            print("shapes:")
            for sp in slide.shapes:
                print(sp.shape_type)
                print(sp.top, sp.left, sp.width, sp.height)

        slide_is_section_header = is_section_header(idx, shapes)
        seen_section_header |= slide_is_section_header
        for shape in shapes:
            if is_title(shape):
                level = 1
                if slide_is_section_header:
                    level = 2
                elif seen_section_header:
                    level = 3

                notes += process_title(shape, idx + 1, level=level)
            elif is_text_block(shape):
                notes += process_text_block(shape, idx + 1)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                notes += process_picture(shape, idx + 1)
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                notes += process_table(shape, idx + 1)

        if not g.disable_notes and slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text
            if text and not has_too_many_links(text):
                notes += process_notes(text, idx + 1)
    out.close()

    if len(notes) > 0:
        print("Process finished with notice:")
        for note in notes:
            print(note)
